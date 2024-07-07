from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from jira import JIRA
import pandas as pd
from modules.utils import get_calendar_week,get_current_year,get_current_month
from modules.jira_operations import get_jira_project_key
import streamlit as st
import tempfile

# create the presentation from a template
def create_powerpoint(jira, jql):
    # Get issues from Jira using JQL
    issues = jira.search_issues(jql)

    # Create DataFrame from issues
    issue_data = []
    for issue in issues:
        issue_data.append({
            'Id': issue.key,
            'Name': issue.fields.summary,
            'Start Date': issue.fields.customfield_10015,
            'Due Date': issue.fields.duedate,
            'Status': issue.fields.status.name,
            'Owner': getattr(issue.fields.assignee, 'displayName', None),
            'Ext.Owner': issue.fields.customfield_10127
        })
    
    df = pd.DataFrame(issue_data)
    st.write(df)
    
    # Path to template and output PowerPoint files (Weekly Report)
    template_path = f'templates/template_{get_jira_project_key()}.pptx' 
    # Load PowerPoint template
    presentation = Presentation(template_path)

    # Copy the second slide (assuming it's the template)
    template_slide = presentation.slides[1]
    slide_layout = template_slide.slide_layout

    # Initialize variables for row count and slide count
    row_count = 0
    slide_count = 1

    # Iterate through rows in the DataFrame with a step of 5
    for start_row in range(0, len(df), 5):
        end_row = start_row + 5
        rows_to_display = df.iloc[start_row:end_row]

        # Create a new slide
        new_slide = presentation.slides.add_slide(slide_layout)
        title_shape = new_slide.shapes.title
        #title_shape.text = f"Actions Overview ({len(issues)} issues)"
        title_shape.text = f"Actions Overview"

        # Add data to the current slide
        add_data_to_slide(new_slide, rows_to_display)

        # Increment slide count
        slide_count += 1

    
    # Update the Calendar Week Placeholders in the slidedeck
    update_current_calendar_week(presentation.slides)
    
    # Construct the final filename by calling the get current Jira Project Key and get Calendar Week
    #presentationFileName = f'{output_path}{get_jira_project_key()}_Weekly_Status_Report_CW_{get_calendar_week()}.pptx'
    
    # Save the presentation
    #Using a with statement is a good practice when dealing with file operations, 
    #as it ensures that the file is properly closed, even if an error occurs during the processing. 
    #In the case of the Presentation class in the python-pptx library, the library itself is designed to handle file operations internally.
    #presentation.save(presentationFileName)

    # Save the presentation to a temporary file
    with tempfile.NamedTemporaryFile(delete=False, suffix=".pptx") as tmpfile:
        presentation.save(tmpfile.name)
        return tmpfile.name
       

# add issues to table
def add_data_to_slide(slide, df):
    # Remove "Content Placeholder 2" shape if it exists
    remove_content_placeholder(slide)

    # Add a table to the slide
    rows, cols = df.shape
    # Specify table dimensions and position
    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(12)
    height = Inches(rows * .5)

    
    table = slide.shapes.add_table(rows=rows + 1, cols=cols, left=left, top=top, width=width, height=height).table
    
    # Add column names to the table header
    for col_num, col_name in enumerate(df.columns):
        table.cell(0, col_num).text = col_name
    
    # Add data to the table
    for row_num, (_, row) in enumerate(df.iterrows(), start=1):
        for col_num, value in enumerate(row):
            table.cell(row_num, col_num).text = str(value)

    format_table(table)
    
# formatting the table
def format_table(table):
    # Adjust column width individually
    column_widths = [Inches(1), Inches(3.5), Inches(1.5), Inches(1.5),Inches(1.5),Inches(1.5),Inches(1.5)]  # Set the desired widths for each column

    for col_num, col_width in enumerate(column_widths):
        table.columns[col_num].width = col_width

    # Adjust font size and font for all cells
    for cell in table.iter_cells():
        cell.text_frame.text = cell.text_frame.text  # This line resets the font settings to default

        # Set font size
        cell.text_frame.paragraphs[0].font.size = Pt(10)  # Set the desired font size, change accordingly

        # Set font (change 'Arial' to your desired font)
        cell.text_frame.paragraphs[0].font.name = 'Open Sauce One'

        # Set text alignment to center
        cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

# remove all placeholders that are auto generated when a new slide is added
def remove_content_placeholder(slide):
    # Iterate through shapes on the slide and remove "Content Placeholder 2" 
    # https://stackoverflow.com/questions/39603318/how-to-delete-unpopulated-placeholder-items-using-python-pptx
    for shape in slide.placeholders:
        if shape.name == "Content Placeholder 2":
            contentPlaceholder = slide.placeholders[1]
            sp = contentPlaceholder.element
            sp.getparent().remove(sp)

def update_current_calendar_week(slides):
    for slide in slides:
        # Iterate through each shape in the current slide
        for shape in slide.shapes:
            #print(shape.shape_type, shape.name)
            if shape.name == "CWPlaceHolder":
                # Replace the text in the shape
                shape.text = str(get_calendar_week())
            if shape.name == "CWPlaceHolderNextSteps":
                # Replace the text in the shape
                shape.text = f'CW {str(get_calendar_week())}'
            if shape.name == "Szöveg helye 3":
                # Replace the text in the shape
                shape.text = f'{get_current_month()} {get_current_year()}'
                

